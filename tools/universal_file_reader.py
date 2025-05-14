"""
Universal File Reader Tool for the AI agent.
Extracts and reads content from various file types including:
- Text files
- Images (OCR)
- PDFs
- Word documents
- PowerPoint presentations
- HTML files
- And more
"""

import sys

import os
import logging
import json
import mimetypes
import traceback
from typing import Dict, List, Optional, Any
from langchain.callbacks.manager import CallbackManagerForToolRun
from langchain.tools.base import BaseTool

# Import specialized libraries for different file types
# Text and HTML extraction
import html2text
import trafilatura

def resource_path(relative_path):
    """
    Get absolute path to resource, works for dev and for PyInstaller
    """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    
    return os.path.join(base_path, relative_path)


# Image processing
try:
    from PIL import Image
    import pytesseract
    IMAGE_SUPPORT = True
except ImportError:
    IMAGE_SUPPORT = False

# PDF support
try:
    import pypdf
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# Document formats
try:
    import docx
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

logger = logging.getLogger(__name__)

class UniversalFileReaderTool(BaseTool):
    """Tool for extracting and reading content from various file types."""
    
    name: str = "universal_file_reader"
    description: str = """
    Extracts and reads content from files of various formats.
    
    Input should be the path to the file.
    Returns the extracted content and metadata, or an error message.
    
    Supported file types include:
    - Text files (.txt, .csv, .json, .xml, etc.)
    - Images (.jpg, .png, .gif, etc.) - uses OCR to extract visible text
    - PDFs (.pdf)
    - Word documents (.docx)
    - PowerPoint presentations (.pptx)
    - HTML files (.html, .htm)
    
    Example: "read_file my_document.pdf" or "extract text from screenshot.jpg"
    """
    
    def _run(self, file_path: str, run_manager: Optional[CallbackManagerForToolRun] = None) -> str:
        """Extract and read content from a file."""
        try:
            # Clean up the input
            file_path = file_path.strip()
            
            # Verify file exists
            if not os.path.exists(file_path):
                return f"Error: File not found: {file_path}"
            
            # Determine file type
            file_type = self._get_file_type(file_path)
            
            # Extract content based on file type
            content = self._extract_content(file_path, file_type)
            
            # Get metadata
            metadata = self._get_metadata(file_path)
            
            # Format response
            response = {
                "success": True,
                "file_path": file_path,
                "file_type": file_type,
                "metadata": metadata,
                "content": content
            }
            
            # Truncate content if it's too long for better readability
            displayed_content = content
            if displayed_content and len(displayed_content) > 1000:
                displayed_content = displayed_content[:997] + "..."
            
            result = f"File: {file_path}\nType: {file_type}\n\nContent:\n{displayed_content}"
            
            # For debugging
            # logger.debug(f"Extracted content from {file_path} ({file_type}): {len(content)} characters")
            
            return result
            
        except Exception as e:
            error_details = traceback.format_exc()
            logger.error(f"Error reading file: {str(e)}\n{error_details}")
            return f"Error reading file: {str(e)}"
    
    def _get_file_type(self, file_path: str) -> str:
        """Determine the file type based on extension and content."""
        # Get file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # Check file extension first
        if ext in ['.txt', '.csv', '.log', '.md', '.json', '.xml', '.yaml', '.yml']:
            return "text"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
            return "image"
        elif ext == '.pdf':
            return "pdf"
        elif ext == '.docx':
            return "docx"
        elif ext == '.pptx':
            return "pptx"
        elif ext in ['.html', '.htm']:
            return "html"
        
        # If extension is not recognized, try to guess from content
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            if mime_type.startswith('text/'):
                return "text"
            elif mime_type.startswith('image/'):
                return "image"
            elif mime_type == 'application/pdf':
                return "pdf"
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return "docx"
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return "pptx"
            elif mime_type in ['text/html', 'application/xhtml+xml']:
                return "html"
        
        # If still not sure, default to text and try to open it
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                f.read(100)  # Try to read the first 100 characters
            return "text"
        except:
            return "binary"  # If we can't read it as text, assume it's binary
    
    def _extract_content(self, file_path: str, file_type: str) -> str:
        """Extract content from a file based on its type."""
        if file_type == "text":
            # Read text files
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error reading text file: {str(e)}")
                return f"Error reading text file: {str(e)}"
                
        elif file_type == "image" and IMAGE_SUPPORT:
            # Use OCR to extract text from images
            try:
                return self._extract_text_from_image(file_path)
            except Exception as e:
                logger.error(f"Error extracting text from image: {str(e)}")
                return f"Error extracting text from image: {str(e)}"
                
        elif file_type == "pdf" and PDF_SUPPORT:
            # Extract text from PDF
            try:
                return self._extract_text_from_pdf(file_path)
            except Exception as e:
                logger.error(f"Error extracting text from PDF: {str(e)}")
                return f"Error extracting text from PDF: {str(e)}"
                
        elif file_type == "docx" and DOCX_SUPPORT:
            # Extract text from Word document
            try:
                return self._extract_text_from_docx(file_path)
            except Exception as e:
                logger.error(f"Error extracting text from Word document: {str(e)}")
                return f"Error extracting text from Word document: {str(e)}"
                
        elif file_type == "pptx" and PPTX_SUPPORT:
            # Extract text from PowerPoint presentation
            try:
                return self._extract_text_from_pptx(file_path)
            except Exception as e:
                logger.error(f"Error extracting text from PowerPoint: {str(e)}")
                return f"Error extracting text from PowerPoint: {str(e)}"
                
        elif file_type == "html":
            # Extract clean text from HTML
            try:
                return self._extract_text_from_html(file_path)
            except Exception as e:
                logger.error(f"Error extracting text from HTML: {str(e)}")
                return f"Error extracting text from HTML: {str(e)}"
                
        else:
            # For unsupported file types or missing dependencies
            supported_types = []
            if IMAGE_SUPPORT:
                supported_types.append("image")
            if PDF_SUPPORT:
                supported_types.append("pdf")
            if DOCX_SUPPORT:
                supported_types.append("docx")
            if PPTX_SUPPORT:
                supported_types.append("pptx")
            
            if file_type in ["image", "pdf", "docx", "pptx"] and file_type not in supported_types:
                return f"Support for {file_type} files is not available. Missing required libraries."
            else:
                return f"Unable to extract content from {file_type} files."
    
    def _extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return text.strip()
        except Exception as e:
            logger.error(f"OCR error: {str(e)}")
            return f"OCR error: {str(e)}"
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        text_content = []
        try:
            with open(file_path, "rb") as pdf_file:
                pdf_reader = pypdf.PdfReader(pdf_file)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text_content.append(page.extract_text())
                
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            return f"PDF extraction error: {str(e)}"
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        try:
            doc = docx.Document(file_path)
            full_text = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                if para.text:
                    full_text.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            return "\n".join(full_text)
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            return f"DOCX extraction error: {str(e)}"
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                slide_content.append(f"\n--- Slide {i+1} ---")
                
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                
                text_content.append("\n".join(slide_content))
            
            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            return f"PPTX extraction error: {str(e)}"
    
    def _extract_text_from_html(self, file_path: str) -> str:
        """Extract clean text from HTML."""
        try:
            # First try with trafilatura for better content extraction
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    html_content = f.read()
                
                extracted_text = trafilatura.extract(html_content)
                if extracted_text:
                    return extracted_text
            except:
                pass  # If trafilatura fails, fall back to html2text
            
            # Fallback to html2text
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            h.ignore_emphasis = True
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            return h.handle(html_content)
        except Exception as e:
            logger.error(f"HTML extraction error: {str(e)}")
            return f"HTML extraction error: {str(e)}"
    
    def _get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get metadata for the file."""
        try:
            stat_info = os.stat(file_path)
            
            metadata = {
                "size_bytes": stat_info.st_size,
                "size_human": self._human_readable_size(stat_info.st_size),
                "created_time": stat_info.st_ctime,
                "modified_time": stat_info.st_mtime,
                "permissions": stat_info.st_mode
            }
            
            return metadata
        except Exception as e:
            logger.error(f"Error getting file metadata: {str(e)}")
            return {"error": str(e)}
    
    def _human_readable_size(self, size_bytes: int) -> str:
        """Convert bytes to human-readable size."""
        units = ['B', 'KB', 'MB', 'GB', 'TB']
        unit_index = 0
        size = float(size_bytes)
        
        while size >= 1024 and unit_index < len(units) - 1:
            size /= 1024
            unit_index += 1
        
        return f"{size:.2f} {units[unit_index]}"